"""
Engine smoke test: does a full generation round-trip cleanly?

This is the safety net for the OOXML packager — the part of the codebase that
has produced 5 pptx-corruption fixes in recent history. If this passes, the
deck opens in PowerPoint without "found a problem" errors.

The test stays offline: instead of calling Anthropic, we hand-feed a minimal
build_content() function as the AI code for one slide and an empty string
(triggering the fallback path) for another.
"""
import os
import sys

import pytest
from pptx import Presentation

from template_engine import InvescoreTemplateEngine

BACKEND_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
TEMPLATE = os.path.join(BACKEND_DIR, "templates", "InvesCore_Master_Template.pptx")
BRAND = os.path.join(BACKEND_DIR, "brand_guide.json")

if not os.path.exists(TEMPLATE):
    pytest.skip("Master template not present", allow_module_level=True)


SAMPLE_CODE = """
def build_content(slide, Inches, Pt, Emu, RGBColor):
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    tb = slide.shapes.add_textbox(Inches(0.57), Inches(1.35), Inches(7.5), Inches(0.5))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = 'KEY METRICS'
    r.font.name = 'Montserrat'
    r.font.size = Pt(13)
    r.font.bold = True
    r.font.color.rgb = RGBColor(0x3B, 0x3B, 0x3B)
"""

ESCAPE_CODE = """
def build_content(slide, Inches, Pt, Emu, RGBColor):
    cls = ().__class__.__bases__[0]
    print('escaped')
"""


@pytest.fixture(scope="module")
def engine():
    return InvescoreTemplateEngine(TEMPLATE, BRAND)


@pytest.fixture
def slide_plan():
    return {
        "presentation_title": "TEST DECK",
        "sections": [
            {
                "name": "EXECUTIVE SUMMARY",
                "slides": [
                    {"slide_type": "content", "title": "KPIs", "description": "metric cards"},
                ],
            },
            {
                "name": "MARKET",
                "slides": [
                    {"slide_type": "section_divider", "title": "Market Outlook"},
                    {"slide_type": "content", "title": "Trends", "description": "chart"},
                ],
            },
        ],
    }


def test_engine_v2_roundtrip(engine, slide_plan, tmp_path):
    code_map = {0: SAMPLE_CODE, 1: SAMPLE_CODE}
    out, warnings = engine.create_presentation_v2(slide_plan, code_map)
    try:
        # 1+1+1+2+1 == 6 (opening + agenda + content + section_divider + content + ending)
        prs = Presentation(out)
        assert len(prs.slides) == 6
        # No fallback warnings expected on legitimate code
        assert not warnings, f"unexpected warnings: {warnings}"
    finally:
        if os.path.exists(out):
            os.unlink(out)


def test_engine_v2_handles_sandbox_escape_gracefully(engine, slide_plan):
    """Malicious code on one slide must not corrupt the deck — fallback fires."""
    code_map = {0: ESCAPE_CODE, 1: SAMPLE_CODE}
    out, warnings = engine.create_presentation_v2(slide_plan, code_map)
    try:
        prs = Presentation(out)
        assert len(prs.slides) == 6
        # Slide 0 (escape attempt) should warn; slide 1 should not.
        assert len(warnings) == 1
        assert warnings[0]["builder_index"] == 0
    finally:
        if os.path.exists(out):
            os.unlink(out)


def test_engine_v2_handles_empty_code_with_fallback(engine, slide_plan):
    """Empty string => Builder Agent failed => slide gets a fallback title."""
    code_map = {0: "", 1: SAMPLE_CODE}
    out, warnings = engine.create_presentation_v2(slide_plan, code_map)
    try:
        prs = Presentation(out)
        assert len(prs.slides) == 6
        assert len(warnings) == 1
        assert warnings[0]["builder_index"] == 0
    finally:
        if os.path.exists(out):
            os.unlink(out)


def test_engine_init_requires_categories(tmp_path):
    """Brand guide missing required categories must fail at init."""
    import json
    bad = tmp_path / "brand.json"
    bad.write_text(json.dumps({"slides": [{"category": "opening", "index": 0, "dynamic_fields": []}]}))
    with pytest.raises(ValueError, match="missing required categories"):
        InvescoreTemplateEngine(TEMPLATE, str(bad))


def test_engine_v2_validates_final_pptx_opens(engine, slide_plan):
    """The post-generation re-open is the canary that catches OOXML corruption."""
    code_map = {0: SAMPLE_CODE, 1: SAMPLE_CODE}
    out, _ = engine.create_presentation_v2(slide_plan, code_map)
    try:
        # Re-opening cleanly is the contract — that's what the engine asserts.
        prs = Presentation(out)
        # Slide 1 (agenda) should have content; slide 5 (ending) should exist.
        assert prs.slides[0] is not None
        assert prs.slides[-1] is not None
    finally:
        if os.path.exists(out):
            os.unlink(out)

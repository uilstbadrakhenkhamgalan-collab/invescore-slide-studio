"""
Sandbox-escape regression tests for the AST-based code validator.

These tests are the contract between the Builder Agent's emitted code and the
sandbox. If any pass that should fail, AI-generated code can escape the sandbox
to subprocess.Popen / file IO / network.
"""
import os

import pytest

from template_engine import InvescoreTemplateEngine


@pytest.fixture(scope="module")
def engine():
    here = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    return InvescoreTemplateEngine(
        os.path.join(here, "templates", "InvesCore_Master_Template.pptx"),
        os.path.join(here, "brand_guide.json"),
    )


# ── Legitimate code paths ────────────────────────────────────────────────────
LEGITIMATE_CODES = [
    pytest.param(
        """
def build_content(slide, Inches, Pt, Emu, RGBColor):
    from pptx.util import Inches, Pt
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    tb.text_frame.text = 'Hello'
""",
        id="basic_textbox",
    ),
    pytest.param(
        """
def build_content(slide, Inches, Pt, Emu, RGBColor):
    import math
    from pptx.util import Inches
    from pptx.dml.color import RGBColor
    n = max(3, math.floor(7 / 2))
    for i in range(n):
        slide.shapes.add_textbox(Inches(1+i), Inches(1), Inches(2), Inches(1))
""",
        id="allowed_math_import",
    ),
    pytest.param(
        """
def build_content(slide, Inches, Pt, Emu, RGBColor):
    from pptx.chart.data import ChartData
    from pptx.enum.chart import XL_CHART_TYPE
    cd = ChartData()
    cd.categories = ['A', 'B', 'C']
    cd.add_series('Vals', (1, 2, 3))
""",
        id="chart_imports",
    ),
]


# ── Sandbox-escape attempts ──────────────────────────────────────────────────
ESCAPE_CODES = [
    pytest.param(
        """
def build_content(slide, Inches, Pt, Emu, RGBColor):
    cls = ().__class__.__bases__[0]
    for c in cls.__subclasses__():
        if 'Popen' in c.__name__:
            c(['ls'])
""",
        id="dunder_walk_to_popen",
    ),
    pytest.param(
        """
def build_content(slide, Inches, Pt, Emu, RGBColor):
    bases = getattr((), '__class__')
""",
        id="getattr_string_dunder",
    ),
    pytest.param(
        """
def build_content(slide, Inches, Pt, Emu, RGBColor):
    target = '__subclasses__'
""",
        id="bare_string_dunder",
    ),
    pytest.param(
        """
def build_content(slide, Inches, Pt, Emu, RGBColor):
    import os
    os.system('ls')
""",
        id="import_os",
    ),
    pytest.param(
        """
def build_content(slide, Inches, Pt, Emu, RGBColor):
    from subprocess import Popen
""",
        id="from_subprocess_import",
    ),
    pytest.param(
        """
def build_content(slide, Inches, Pt, Emu, RGBColor):
    exec('print(1)')
""",
        id="exec_call",
    ),
    pytest.param(
        """
def build_content(slide, Inches, Pt, Emu, RGBColor):
    eval('1+1')
""",
        id="eval_call",
    ),
    pytest.param(
        """
def build_content(slide, Inches, Pt, Emu, RGBColor):
    open('/etc/passwd').read()
""",
        id="open_call",
    ),
    pytest.param(
        """
def build_content(slide, Inches, Pt, Emu, RGBColor):
    g = globals()
""",
        id="globals_call",
    ),
    pytest.param(
        """
def build_content(slide, Inches, Pt, Emu, RGBColor):
    b = __builtins__
""",
        id="builtins_reference",
    ),
    pytest.param(
        """
def build_content(slide, Inches, Pt, Emu, RGBColor):
    f = (lambda: None).__code__
""",
        id="code_attribute_walk",
    ),
    pytest.param(
        """
def build_content(slide, Inches, Pt, Emu, RGBColor):
    setattr(slide, '__class__', None)
""",
        id="setattr_dunder_attr",
    ),
    pytest.param(
        "def build_content(:):",
        id="syntax_error",
    ),
    pytest.param(
        """
def something_else():
    pass
""",
        id="missing_build_content",
    ),
    pytest.param(
        """
def build_content(slide, Inches, Pt, Emu, RGBColor):
    pass
def build_content(slide, Inches, Pt, Emu, RGBColor):
    pass
""",
        id="duplicate_build_content",
    ),
    pytest.param(
        """
def build_content(slide, Inches, Pt, Emu, RGBColor):
    import importlib
    importlib.import_module('os')
""",
        id="importlib_bypass",
    ),
]


@pytest.mark.parametrize("code", LEGITIMATE_CODES)
def test_legitimate_code_passes_validation(engine, code):
    ok, reason = engine._validate_code(code)
    assert ok, f"legitimate code rejected: {reason}"


@pytest.mark.parametrize("code", ESCAPE_CODES)
def test_escape_attempts_blocked(engine, code):
    ok, reason = engine._validate_code(code)
    assert not ok, "escape attempt should have been rejected"
    assert reason, "rejection should include a reason"

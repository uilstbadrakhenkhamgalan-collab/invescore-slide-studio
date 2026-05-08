"""
Concurrency regression: multiple simultaneous generations must not corrupt
each other's tmp/ artifacts (filesystem-state hazard) and the cleanup sweep
must not delete in-flight artifacts.
"""
import asyncio
import json
import os
import time

import pytest
from pptx import Presentation

from template_engine import InvescoreTemplateEngine
import main as backend_main

BACKEND_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
TEMPLATE = os.path.join(BACKEND_DIR, "templates", "InvesCore_Master_Template.pptx")
BRAND = os.path.join(BACKEND_DIR, "brand_guide.json")

if not os.path.exists(TEMPLATE):
    pytest.skip("Master template not present", allow_module_level=True)

SAMPLE = """
def build_content(slide, Inches, Pt, Emu, RGBColor):
    from pptx.util import Inches, Pt
    tb = slide.shapes.add_textbox(Inches(0.57), Inches(1.35), Inches(7), Inches(0.5))
    tb.text_frame.text = 'X'
"""


def _plan(title: str):
    return {
        "presentation_title": title,
        "sections": [
            {
                "name": "S",
                "slides": [
                    {"slide_type": "content", "title": "T", "description": "d"},
                ],
            }
        ],
    }


@pytest.mark.asyncio
async def test_concurrent_generations_do_not_collide():
    """Three engines running in parallel must each produce a valid deck."""
    engine = InvescoreTemplateEngine(TEMPLATE, BRAND)

    async def run(name: str):
        out, _ = await asyncio.to_thread(
            engine.create_presentation_v2, _plan(name), {0: SAMPLE}
        )
        return out

    paths = await asyncio.gather(*(run(f"deck-{i}") for i in range(3)))
    try:
        for p in paths:
            assert os.path.exists(p)
            prs = Presentation(p)
            assert len(prs.slides) >= 4
        # Each generation must produce a different temp path
        assert len(set(paths)) == 3
    finally:
        for p in paths:
            if os.path.exists(p):
                os.unlink(p)


def test_cleanup_skips_in_flight_artifacts():
    """An artifact written within the last 60s must not be reaped."""
    backend_main.TMP_DIR.mkdir(parents=True, exist_ok=True)
    artifact_id = "f" * 32
    meta = backend_main._artifact_meta_path(artifact_id)
    pptx = backend_main._artifact_file_path(artifact_id)
    try:
        # Brand new, just-created — but already "expired" — should still survive.
        meta.write_text(json.dumps({
            "filename": "x.pptx",
            "token_sha256": backend_main._hash_download_token("t"),
            "created_at": time.time(),  # just now
            "expires_at": time.time() - 1,  # already expired
        }))
        pptx.write_bytes(b"PK")
        backend_main._cleanup_expired_artifacts()
        assert pptx.exists(), "in-flight artifact was reaped before grace period"
        assert meta.exists(), "in-flight metadata was reaped before grace period"
    finally:
        meta.unlink(missing_ok=True)
        pptx.unlink(missing_ok=True)


def test_cleanup_removes_old_expired_artifacts():
    backend_main.TMP_DIR.mkdir(parents=True, exist_ok=True)
    artifact_id = "e" * 32
    meta = backend_main._artifact_meta_path(artifact_id)
    pptx = backend_main._artifact_file_path(artifact_id)
    try:
        meta.write_text(json.dumps({
            "filename": "old.pptx",
            "token_sha256": backend_main._hash_download_token("t"),
            "created_at": time.time() - 3600,  # an hour old
            "expires_at": time.time() - 100,
        }))
        pptx.write_bytes(b"PK")
        removed = backend_main._cleanup_expired_artifacts()
        assert not pptx.exists()
        assert not meta.exists()
        assert removed >= 1
    finally:
        meta.unlink(missing_ok=True)
        pptx.unlink(missing_ok=True)

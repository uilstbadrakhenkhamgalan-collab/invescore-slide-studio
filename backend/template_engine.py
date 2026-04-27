"""
InvesCore Slide Studio â€” Template Engine v3
Hybrid approach:
  - Clone-based for opening / agenda / section_divider / closing
  - Blank-slide + brand-frame + AI-generated python-pptx for content slides
"""
import json
import os
import copy
import shutil
import subprocess
import sys
import tempfile
import threading
import traceback
import zipfile
import re
from pathlib import Path
from lxml import etree
from pptx import Presentation
from pptx.util import Pt, Inches, Emu
from pptx.oxml.ns import qn
from pptx.dml.color import RGBColor


# â”€â”€ XML namespaces â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
NS_REL = "http://schemas.openxmlformats.org/package/2006/relationships"
NS_PML = "http://schemas.openxmlformats.org/presentationml/2006/main"
NS_RID = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


# â”€â”€ Brand-frame constants (derived from content_text slide, index 15) â”€â”€â”€â”€â”€â”€â”€â”€â”€

# Shape IDs in the content_text slide that belong to the CONTENT AREA.
# These are stripped when repurposing the cloned slide as a brand-frame base.
CONTENT_AREA_SHAPE_IDS = frozenset({2, 6, 7, 8, 9, 10, 11, 12, 13, 14, 15})

# Nav-bar section-label shape IDs in left-to-right display order (8 slots).
NAV_LABEL_SHAPE_IDS = [21, 22, 23, 24, 25, 27, 28, 29]

# Shape ID of the large section-name label below the nav bar.
CONTENT_SECTION_TITLE_SHAPE_ID = 18

# Shape ID of the page-number text at bottom-right.
CONTENT_PAGE_NUM_SHAPE_ID = 30

# Agenda slide: section-title shape IDs for slots 1-8 (left col first, then right col).
AGENDA_SECTION_TITLE_IDS = [9, 15, 21, 27, 10, 16, 22, 28]

# Agenda slide: page-range shape IDs for slots 1-8.
AGENDA_SECTION_PAGES_IDS = [31, 32, 33, 34, 35, 36, 37, 38]

# Usable content-area bounds in inches (slide is 10" Ã— 5.625").
# Builder Agent must keep all content within these limits.
CONTENT_AREA_BOUNDS = {
    "top":    1.30,   # below header bar (0.59") + section-title shape (0.95")
    "bottom": 5.26,   # above page-number area
    "left":   0.57,   # left margin
    "right":  9.28,   # before right-side decorative elements
    # Working area: ~8.71" wide Ã— ~3.96" tall
}


# â”€â”€ Code-execution safety â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€

_BLOCKED_IMPORTS = frozenset([
    "os", "sys", "subprocess", "requests", "socket", "shutil", "pathlib",
    "io", "http", "urllib", "tempfile", "glob", "importlib", "builtins",
    "pickle", "ctypes", "multiprocessing", "threading", "ftplib", "smtplib",
    "zipfile", "tarfile", "sqlite3", "dbm", "shelve", "signal", "pty",
])

_BLOCKED_TOKENS = [
    "exec(",    "eval(",    "__import__(", "open(",    "compile(",
    "__builtins__", "globals(", "locals(",  "vars(",   "breakpoint(",
]

# Only these top-level module families are allowed inside build_content().
_ALLOWED_MODULES = frozenset([
    "pptx", "math", "datetime", "random", "itertools",
    "collections", "decimal", "fractions", "statistics", "functools",
])

_EXEC_TIMEOUT_SEC = 10
_WORKER_TIMEOUT_SEC = 20
_WORKER_SCRIPT_PATH = Path(__file__).with_name("slide_builder_worker.py")


# â”€â”€ Engine â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€

class InvescoreTemplateEngine:

    def __init__(self, template_path: str, brand_guide_path: str):
        self.template_path = template_path
        self.brand_guide_path = brand_guide_path
        with open(brand_guide_path, encoding="utf-8") as f:
            self.brand = json.load(f)
        self.category_map = {s["category"]: s for s in self.brand["slides"]}

    # â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•
    # V1 Public API  (unchanged â€” used by existing /api/generate endpoint)
    # â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•

    def create_presentation(self, slide_plan: list[dict]) -> str:
        """
        V1: Clone-based, text-swap only.
        slide_plan: [{"template": "opening", "content": {...}}, ...]
        Returns path to output .pptx (caller responsible for deletion).
        """
        source_indices = []
        for spec in slide_plan:
            cat = spec["template"]
            if cat not in self.category_map:
                raise ValueError(
                    f"Unknown template: '{cat}'. Available: {list(self.category_map)}"
                )
            source_indices.append(self.category_map[cat]["index"])

        output_path = self._build_pptx_via_zip(source_indices)

        prs = Presentation(output_path)
        for spec, slide in zip(slide_plan, prs.slides):
            self._apply_content(slide, spec["template"], spec.get("content", {}))
        prs.save(output_path)

        return output_path

    def create_presentation_to_file(self, slide_plan: list[dict], output_path: str):
        """V1: Create presentation and save to a specific path."""
        tmp = self.create_presentation(slide_plan)
        shutil.move(tmp, output_path)

    # â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•
    # V2 Public API  (new â€” hybrid clone + AI-content)
    # â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•

    def create_presentation_v2(
        self, slide_plan: dict, content_code_map: dict
    ) -> tuple[str, list[dict]]:
        """
        V2: Hybrid â€” clone for structural slides, AI python-pptx for content.

        Args:
            slide_plan:
                {
                  "presentation_title": str,
                  "sections": [
                    {
                      "name": str,
                      "slides": [
                        {
                          "slide_type": "content" | "section_divider",
                          "title": str,
                          "description": str,
                          "content_spec": {...}   # for content slides
                        }
                      ]
                    }
                  ]
                }
            content_code_map:
                { content_slide_index (0-based): python_code_string }
                Index counts only content slides, not
                opening / agenda / closing.

        Returns:
            Tuple of generated .pptx path and non-fatal content-render warnings.
        """
        all_sections = [s["name"] for s in slide_plan["sections"]]

        source_indices = [
            self.category_map["opening"]["index"],
            self.category_map["agenda"]["index"],
        ]

        slide_meta = []
        builder_idx = 0
        pptx_idx = 2  # 0=opening, 1=agenda

        for section in slide_plan["sections"]:
            for slide_spec in section.get("slides", []):
                stype = slide_spec.get("slide_type", "content")
                if stype == "section_divider":
                    source_indices.append(
                        self.category_map["section_divider"]["index"]
                    )
                    current_builder_idx = None
                else:
                    source_indices.append(
                        self.category_map["content_text"]["index"]
                    )
                    current_builder_idx = builder_idx
                    builder_idx += 1

                slide_meta.append({
                    "pptx_idx": pptx_idx,
                    "section_name": section["name"],
                    "slide_spec": slide_spec,
                    "stype": stype,
                    "builder_idx": current_builder_idx,
                })
                pptx_idx += 1

        source_indices.append(self.category_map["ending"]["index"])

        output_path = self._build_pptx_via_zip(source_indices)
        prs = Presentation(output_path)

        self._apply_content(
            prs.slides[0], "opening",
            {"presentation_title": slide_plan["presentation_title"]}
        )
        self._apply_agenda_v2(prs.slides[1], slide_plan["sections"], start_page=3)

        last_slide_label = "unknown"
        for meta in slide_meta:
            if meta["stype"] != "section_divider":
                continue

            pidx = meta["pptx_idx"]
            section_name = meta["section_name"]
            slide_spec = meta["slide_spec"]
            slide = prs.slides[pidx]
            page_number = pidx + 1
            slide_title = slide_spec.get("title", f"Slide {page_number}")
            slide_label = f"slide={page_number} section='{section_name}' title='{slide_title}'"
            last_slide_label = slide_label

            print(f"[engine] Preparing {slide_label} type=section_divider")
            self._apply_content(slide, "section_divider", {
                "section_title":       slide_spec.get("title", section_name),
                "section_description": "",  # description is internal AI guidance only
            })
            print(f"[engine] Section divider applied â€” {slide_label}")

        print(f"[engine] Saving structural slides to {output_path}")
        prs.save(output_path)

        warnings: list[dict] = []

        for meta in slide_meta:
            if meta["stype"] == "section_divider":
                continue

            pidx = meta["pptx_idx"]
            section_name = meta["section_name"]
            slide_spec = meta["slide_spec"]
            current_builder_idx = meta["builder_idx"]
            page_number = pidx + 1
            slide_title = slide_spec.get("title", f"Slide {page_number}")
            slide_label = f"slide={page_number} section='{section_name}' title='{slide_title}'"
            last_slide_label = slide_label
            code = content_code_map.get(current_builder_idx, "")

            print(
                f"[engine] Rendering {slide_label} "
                f"builder_idx={current_builder_idx}"
            )

            warning = self._apply_content_slide_with_worker(
                presentation_path=output_path,
                slide_index=pidx,
                section_name=section_name,
                all_sections=all_sections,
                page_number=page_number,
                slide_title=slide_title,
                code=code,
                slide_label=slide_label,
            )
            if warning:
                warnings.append({
                    "builder_index": current_builder_idx,
                    "page_number": page_number,
                    "section_name": section_name,
                    "title": slide_title,
                    "message": warning,
                })

        print(f"[engine] Validating saved .pptx...")
        try:
            test_prs = Presentation(output_path)
            slide_count = len(test_prs.slides)
            print(f"[engine] Validation OK â€” {slide_count} slides readable")
        except Exception as val_exc:
            print(f"[engine] VALIDATION FAILED â€” last slide was: {last_slide_label}")
            print(f"[engine] Validation error: {val_exc}")
            raise RuntimeError(
                f"Generated .pptx is corrupt (failed to re-open): {val_exc}. "
                f"Last slide processed: {last_slide_label}"
            )

        return output_path, warnings


    def get_content_area_bounds(self) -> dict:
        """
        Returns the usable content-area boundaries (inches) for Builder Agent use.
        All generated shapes must stay within these limits.
        """
        return dict(CONTENT_AREA_BOUNDS)

    # â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•
    # Brand-frame manipulation  (v2 helpers)
    # â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•

    def _clear_content_area(self, slide):
        """
        Remove content-area shapes from a cloned content_text slide,
        leaving only brand-frame elements (header bar, logo, nav, page numberâ€¦).
        """
        sp_tree = slide.shapes._spTree
        to_remove = [
            shape._element
            for shape in slide.shapes
            if shape.shape_id in CONTENT_AREA_SHAPE_IDS
        ]
        for elem in to_remove:
            sp_tree.remove(elem)

    def _update_brand_frame(
        self, slide, active_section: str, all_sections: list, page_number: int
    ):
        """
        Refresh the dynamic brand-frame elements for a specific slide:
          - Nav-bar labels (up to 8): update text, bold the active section
          - Section-title label below nav bar
          - Page number
        """
        try:
            active_idx = all_sections.index(active_section)
        except ValueError:
            active_idx = 0

        for shape in slide.shapes:
            sid = shape.shape_id

            # â”€â”€ Nav section labels â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
            if sid in NAV_LABEL_SHAPE_IDS:
                slot = NAV_LABEL_SHAPE_IDS.index(sid)
                if not shape.has_text_frame:
                    continue
                new_text = all_sections[slot] if slot < len(all_sections) else ""
                is_active = slot == active_idx
                tf = shape.text_frame
                if tf.paragraphs and tf.paragraphs[0].runs:
                    run = tf.paragraphs[0].runs[0]
                    run.text = new_text
                    run.font.bold = is_active
                    # Wipe any extra runs in the paragraph
                    for extra in tf.paragraphs[0].runs[1:]:
                        extra.text = ""
                elif tf.paragraphs and new_text:
                    p_elem = tf.paragraphs[0]._p
                    r = etree.SubElement(p_elem, qn("a:r"))
                    t = etree.SubElement(r, qn("a:t"))
                    t.text = new_text

            # â”€â”€ Section title below nav bar â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
            elif sid == CONTENT_SECTION_TITLE_SHAPE_ID:
                if shape.has_text_frame:
                    self._set_shape_text(shape, active_section)

            # â”€â”€ Page number â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
            elif sid == CONTENT_PAGE_NUM_SHAPE_ID:
                if shape.has_text_frame:
                    self._set_shape_text(shape, str(page_number))

        # Add rotated "Confidential document" label at the right edge
        self._ensure_confidential_label(slide)

    def _ensure_confidential_label(self, slide):
        """Add a rotated 'Confidential document' text box at the right edge if absent."""
        CONFIDENTIAL_TEXT = "Confidential document"
        # Check if label already exists (avoid duplicates on repeated calls)
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip() == CONFIDENTIAL_TEXT:
                return
        from pptx.util import Pt, Inches, Emu
        from pptx.dml.color import RGBColor
        # Position: right edge of slide, spanning full slide height
        # Slide dimensions: 10" x 5.63" (standard widescreen)
        txBox = slide.shapes.add_textbox(
            Inches(9.55), Inches(0.1), Inches(0.35), Inches(5.43)
        )
        tf = txBox.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = 2  # PP_ALIGN.CENTER
        run = p.add_run()
        run.text = CONFIDENTIAL_TEXT
        run.font.size = Pt(7)
        run.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
        run.font.bold = False
        # Rotate 90 degrees counter-clockwise (text reads bottom-to-top)
        txBox.rotation = 90

    def _apply_agenda_v2(self, slide, sections: list, start_page: int):
        """
        Fill the agenda slide for up to 8 sections, calculating page ranges.

        sections:   [{"name": str, "slides": [...]}, ...]
        start_page: 1-based page number of the first content slide.
        """
        # Compute page ranges
        page_ranges: list[str] = []
        cur = start_page
        for sec in sections:
            count = len(sec.get("slides", []))
            if count == 0:
                page_ranges.append("")
            elif count == 1:
                page_ranges.append(f"pg. {cur}")
            else:
                page_ranges.append(f"pg. {cur}â€“{cur + count - 1}")
            cur += count

        for shape in slide.shapes:
            sid = shape.shape_id

            if sid in AGENDA_SECTION_TITLE_IDS:
                slot = AGENDA_SECTION_TITLE_IDS.index(sid)
                new_text = sections[slot]["name"] if slot < len(sections) else ""
                if shape.has_text_frame:
                    self._set_shape_text(shape, new_text)

            elif sid in AGENDA_SECTION_PAGES_IDS:
                slot = AGENDA_SECTION_PAGES_IDS.index(sid)
                new_text = (
                    page_ranges[slot] if slot < len(page_ranges) else ""
                )
                if shape.has_text_frame:
                    self._set_shape_text(shape, new_text)

    def _set_shape_text(self, shape, text: str):
        """
        Replace the text in a shape's first paragraph / first run,
        preserving all existing font formatting.
        """
        tf = shape.text_frame
        if not tf.paragraphs:
            return
        para = tf.paragraphs[0]
        if para.runs:
            para.runs[0].text = text
            for extra in para.runs[1:]:
                extra.text = ""
        else:
            p_elem = para._p
            r = etree.SubElement(p_elem, qn("a:r"))
            t = etree.SubElement(r, qn("a:t"))
            t.text = text

    def _add_fallback_title(self, slide, title: str):
        """
        Add a plain title text box when AI code execution fails,
        so the slide is not completely blank.
        """
        if not title:
            return
        b = CONTENT_AREA_BOUNDS
        left   = Inches(b["left"])
        top    = Inches(b["top"])
        width  = Inches(b["right"] - b["left"])
        height = Inches(0.6)

        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title.upper()
        run.font.name = "Montserrat"
        run.font.size = Pt(18)
        run.font.color.rgb = RGBColor(0x3B, 0x3B, 0x3B)

    def _prepare_content_slide(
        self, slide, section_name: str, all_sections: list[str], page_number: int
    ):
        """Reset a cloned content slide back to its clean brand-frame base."""
        self._clear_content_area(slide)
        self._update_brand_frame(slide, section_name, all_sections, page_number)

    def _apply_ai_content_to_slide(
        self,
        slide,
        *,
        section_name: str,
        all_sections: list[str],
        page_number: int,
        slide_title: str,
        code: str,
        slide_label: str,
    ) -> str | None:
        """
        Populate a single content slide and return a warning message when the
        builder code had to fall back to a plain title.
        """
        self._prepare_content_slide(slide, section_name, all_sections, page_number)

        if code:
            success = self._execute_content_code(slide, code, slide_label)
            if success:
                print(f"[engine] Content built successfully for {slide_label}")
                return None

            print(f"[engine] Using fallback title for {slide_label}")
            self._add_fallback_title(slide, slide_title)
            return (
                "AI content rendering failed validation or execution. "
                "A fallback title slide was used."
            )

        print(f"[engine] No builder code available â€” using fallback title for {slide_label}")
        self._add_fallback_title(slide, slide_title)
        return "No builder code was available for this slide. A fallback title slide was used."

    def _apply_fallback_content_to_file(
        self,
        *,
        presentation_path: str,
        slide_index: int,
        section_name: str,
        all_sections: list[str],
        page_number: int,
        slide_title: str,
        slide_label: str,
    ):
        """Safely write a fallback title slide when the worker fails outright."""
        prs = Presentation(presentation_path)
        slide = prs.slides[slide_index]
        self._prepare_content_slide(slide, section_name, all_sections, page_number)
        self._add_fallback_title(slide, slide_title)
        prs.save(presentation_path)
        print(f"[engine] Worker failed â€” fallback title saved for {slide_label}")

    def _apply_content_slide_with_worker(
        self,
        *,
        presentation_path: str,
        slide_index: int,
        section_name: str,
        all_sections: list[str],
        page_number: int,
        slide_title: str,
        code: str,
        slide_label: str,
    ) -> str | None:
        """
        Render a content slide in an isolated Python worker process and atomically
        replace the deck only when the worker succeeds.
        """
        task_path = ""
        output_copy = ""
        task = {
            "template_path": self.template_path,
            "brand_guide_path": self.brand_guide_path,
            "source_path": presentation_path,
            "slide_index": slide_index,
            "section_name": section_name,
            "all_sections": all_sections,
            "page_number": page_number,
            "slide_title": slide_title,
            "slide_label": slide_label,
            "code": code,
        }

        try:
            with tempfile.NamedTemporaryFile(
                mode="w", suffix=".pptx", delete=False
            ) as output_file:
                output_copy = output_file.name
            task["output_path"] = output_copy

            with tempfile.NamedTemporaryFile(
                mode="w", suffix=".json", delete=False, encoding="utf-8"
            ) as task_file:
                json.dump(task, task_file, ensure_ascii=False)
                task_path = task_file.name

            result = subprocess.run(
                [sys.executable, str(_WORKER_SCRIPT_PATH), task_path],
                capture_output=True,
                text=True,
                timeout=_WORKER_TIMEOUT_SEC,
                check=False,
            )

            if result.returncode != 0:
                stderr = result.stderr.strip() or result.stdout.strip()
                raise RuntimeError(stderr or "worker exited with a non-zero status")

            stdout = result.stdout.strip()
            lines = stdout.splitlines()
            payload = json.loads(lines[-1] if lines else "{}")

            if not output_copy or not os.path.exists(output_copy):
                raise RuntimeError("worker finished without producing an output deck")

            os.replace(output_copy, presentation_path)
            output_copy = ""
            return payload.get("warning")

        except subprocess.TimeoutExpired:
            self._apply_fallback_content_to_file(
                presentation_path=presentation_path,
                slide_index=slide_index,
                section_name=section_name,
                all_sections=all_sections,
                page_number=page_number,
                slide_title=slide_title,
                slide_label=slide_label,
            )
            return "Isolated content worker timed out. A fallback title slide was used."

        except Exception as exc:
            print(f"[engine] Worker failed for {slide_label}: {exc}")
            self._apply_fallback_content_to_file(
                presentation_path=presentation_path,
                slide_index=slide_index,
                section_name=section_name,
                all_sections=all_sections,
                page_number=page_number,
                slide_title=slide_title,
                slide_label=slide_label,
            )
            return (
                "Isolated content worker failed before rendering the slide. "
                "A fallback title slide was used."
            )

        finally:
            if task_path and os.path.exists(task_path):
                os.unlink(task_path)
            if output_copy and os.path.exists(output_copy):
                os.unlink(output_copy)

    # â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•
    # Code execution sandbox
    # â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•

    def _validate_code(self, code: str) -> tuple[bool, str]:
        """
        Static analysis of AI-generated code for dangerous patterns.
        Returns (is_safe, reason_if_rejected).
        """
        for mod in _BLOCKED_IMPORTS:
            if (
                f"import {mod}" in code
                or f"from {mod} " in code
                or f"from {mod}." in code
            ):
                return False, f"blocked import: {mod}"
        for token in _BLOCKED_TOKENS:
            if token in code:
                return False, f"blocked token: {token!r}"
        return True, ""

    def _execute_content_code(self, slide, code_string: str, slide_label: str = "") -> bool:
        """
        Execute Builder Agent's python-pptx code inside a safety sandbox.

        The code must define:
            def build_content(slide, Inches, Pt, Emu, RGBColor): ...

        Safety layers:
        1. Static validation â€” blocked imports / tokens
        2. Restricted __import__ â€” only pptx, math, datetime, etc.
        3. Restricted builtins â€” no open(), exec(), eval(), globals()â€¦
        4. Threading timeout â€” 10 s hard limit per slide
        5. Rollback â€” any shapes added before a crash are removed
        6. XML sanity check â€” slide XML must be serialisable after exec

        Returns True on success, False on any rejection / error / timeout.
        """
        prefix = f"[builder{' ' + slide_label if slide_label else ''}]"

        # â”€â”€ Static validation â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
        is_safe, reason = self._validate_code(code_string)
        if not is_safe:
            print(f"{prefix} VALIDATION FAILED â€” {reason}")
            print(f"{prefix} Rejected code (first 500 chars): {code_string[:500]}")
            return False

        # â”€â”€ Snapshot existing shape elements for rollback â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
        sp_tree = slide.shapes._spTree
        snapshot = copy.deepcopy(sp_tree)

        def _restore_snapshot():
            for child in list(sp_tree):
                sp_tree.remove(child)
            for child in list(snapshot):
                sp_tree.append(copy.deepcopy(child))

        result     = [False]
        exc_holder = [None]

        def _safe_import(name, *args, **kwargs):
            top = name.split(".")[0]
            if top not in _ALLOWED_MODULES:
                raise ImportError(
                    f"Module '{name}' is not allowed inside build_content"
                )
            return __import__(name, *args, **kwargs)

        def _run():
            try:
                safe_builtins = {
                    # Arithmetic / iteration
                    "range": range, "len": len, "str": str, "int": int,
                    "float": float, "list": list, "dict": dict, "tuple": tuple,
                    "bool": bool, "set": set, "frozenset": frozenset,
                    "enumerate": enumerate, "zip": zip, "map": map,
                    "filter": filter, "round": round, "max": max, "min": min,
                    "abs": abs, "sum": sum, "sorted": sorted,
                    "reversed": reversed, "any": any, "all": all,
                    # Light inspection helpers for defensive code paths
                    "isinstance": isinstance,
                    "getattr": getattr,
                    # Constants
                    "True": True, "False": False, "None": None,
                    # Controlled import gateway
                    "__import__": _safe_import,
                }
                namespace = {"__builtins__": safe_builtins}
                exec(code_string, namespace)   # noqa: S102
                build_fn = namespace.get("build_content")
                if build_fn is None:
                    print(f"{prefix} No build_content function found in code")
                    return
                build_fn(slide, Inches, Pt, Emu, RGBColor)
                result[0] = True
            except Exception as exc:
                exc_holder[0] = exc

        thread = threading.Thread(target=_run, daemon=True)
        thread.start()
        thread.join(timeout=_EXEC_TIMEOUT_SEC)

        if thread.is_alive():
            print(f"{prefix} TIMED OUT after {_EXEC_TIMEOUT_SEC}s â€” rolling back and using fallback title")
            _restore_snapshot()
            return False

        if exc_holder[0]:
            exc = exc_holder[0]
            print(f"{prefix} EXECUTION FAILED:")
            print(f"{prefix}   {type(exc).__name__}: {exc}")
            print(f"{prefix}   Traceback:")
            tb_lines = traceback.format_exception(type(exc), exc, exc.__traceback__)
            for line in tb_lines:
                for subline in line.splitlines():
                    print(f"{prefix}     {subline}")
            _restore_snapshot()
            print(f"{prefix}   Rolled back to pre-execution snapshot")
            return False

        # â”€â”€ XML sanity check â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
        try:
            etree.tostring(slide._element)
            print(f"{prefix} Execution OK â€” XML sanity check passed")
        except Exception as xml_exc:
            print(f"{prefix} XML SANITY CHECK FAILED after execution: {xml_exc}")
            _restore_snapshot()
            return False

        return result[0]

    # â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•
    # ZIP-level slide builder  (unchanged from v2 â€” handles clone logic)
    # â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•

    def _build_pptx_via_zip(self, source_indices: list[int]) -> str:
        """
        Core engine: opens the template ZIP, duplicates slides as needed,
        produces a new .pptx with exactly the requested slides in order.
        Returns path to temp output file.
        """
        with zipfile.ZipFile(self.template_path, "r") as src_zip:
            # 1. Parse presentation.xml to get the slide rId list
            prs_xml_bytes  = src_zip.read("ppt/presentation.xml")
            prs_xml        = etree.fromstring(prs_xml_bytes)

            prs_rels_bytes = src_zip.read("ppt/_rels/presentation.xml.rels")
            prs_rels       = etree.fromstring(prs_rels_bytes)

            # Map rId â†’ Target (e.g. 'slides/slide1.xml')
            rId_to_target = {}
            for rel in prs_rels.findall(f"{{{NS_REL}}}Relationship"):
                rId_to_target[rel.get("Id")] = rel.get("Target")

            # Get ordered slide rIds from presentation.xml
            sldIdLst = prs_xml.find(f".//{{{NS_PML}}}sldIdLst")
            if sldIdLst is None:
                raise ValueError(
                    "Template PPTX is missing sldIdLst — the file may be corrupt"
                )
            original_slide_rids = [
                sldId.get(f"{{{NS_RID}}}id")
                for sldId in sldIdLst.findall(f"{{{NS_PML}}}sldId")
            ]

            def target_to_zip_path(target):
                return "ppt/" + target.lstrip("/")

            original_slide_zippaths = [
                target_to_zip_path(rId_to_target[rid])
                for rid in original_slide_rids
            ]

            # 2. Resolve desired source zip-paths
            desired_zippaths = [original_slide_zippaths[i] for i in source_indices]

            # 3. Build new ZIP
            with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp_f:
                tmp_path = tmp_f.name

            slide_pattern      = re.compile(r"^ppt/slides/slide\d+\.xml$")
            slide_rels_pattern = re.compile(r"^ppt/slides/_rels/slide\d+\.xml\.rels$")
            notes_slide_pattern      = re.compile(r"^ppt/notesSlides/notesSlide\d+\.xml$")
            notes_slide_rels_pattern = re.compile(r"^ppt/notesSlides/_rels/notesSlide\d+\.xml\.rels$")
            NOTES_REL_TYPE = (
                "http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesSlide"
            )

            # Parse [Content_Types].xml so we can rebuild it accurately
            NS_CT = "http://schemas.openxmlformats.org/package/2006/content-types"
            ct_xml = etree.fromstring(src_zip.read("[Content_Types].xml"))

            with zipfile.ZipFile(tmp_path, "w", zipfile.ZIP_DEFLATED) as dst_zip:
                # Copy all non-slide, non-content-types files verbatim
                for name in src_zip.namelist():
                    if name in (
                        "ppt/presentation.xml",
                        "ppt/_rels/presentation.xml.rels",
                        "[Content_Types].xml",
                    ):
                        continue
                    if slide_pattern.match(name) or slide_rels_pattern.match(name):
                        continue
                    if notes_slide_pattern.match(name) or notes_slide_rels_pattern.match(name):
                        continue
                    dst_zip.writestr(name, src_zip.read(name))

                # Write desired slides with sequential new names
                new_slide_rids    = []
                new_slide_targets = []

                for i, src_zippath in enumerate(desired_zippaths):
                    new_num      = i + 1
                    new_name     = f"ppt/slides/slide{new_num}.xml"
                    new_rels_name = f"ppt/slides/_rels/slide{new_num}.xml.rels"
                    new_target   = f"slides/slide{new_num}.xml"

                    dst_zip.writestr(new_name, src_zip.read(src_zippath))

                    src_rels_path = src_zippath.replace(
                        "ppt/slides/", "ppt/slides/_rels/"
                    ) + ".rels"
                    if src_rels_path in src_zip.namelist():
                        rels_xml = etree.fromstring(src_zip.read(src_rels_path))
                        for rel in list(rels_xml.findall(f"{{{NS_REL}}}Relationship")):
                            if rel.get("Type") == NOTES_REL_TYPE:
                                rels_xml.remove(rel)
                        dst_zip.writestr(new_rels_name, etree.tostring(
                            rels_xml, xml_declaration=True, encoding="UTF-8", standalone=True,
                        ))

                    new_slide_targets.append(new_target)
                    new_slide_rids.append(f"rId_slide{new_num}")

                # Write updated [Content_Types].xml â€” remove stale slide Override
                # entries and add exact entries for the slides we actually generated.
                SLIDE_CT = (
                    "application/vnd.openxmlformats-officedocument"
                    ".presentationml.slide+xml"
                )
                slide_part_re = re.compile(r"^/ppt/slides/slide\d+\.xml$")
                notes_part_re = re.compile(r"^/ppt/notesSlides/notesSlide\d+\.xml$")
                for override in list(ct_xml.findall(f"{{{NS_CT}}}Override")):
                    part_name = override.get("PartName", "")
                    if slide_part_re.match(part_name) or notes_part_re.match(part_name):
                        ct_xml.remove(override)
                for i in range(len(desired_zippaths)):
                    elem = etree.SubElement(ct_xml, f"{{{NS_CT}}}Override")
                    elem.set("PartName", f"/ppt/slides/slide{i + 1}.xml")
                    elem.set("ContentType", SLIDE_CT)
                dst_zip.writestr(
                    "[Content_Types].xml",
                    etree.tostring(
                        ct_xml,
                        xml_declaration=True,
                        encoding="UTF-8",
                        standalone=True,
                    ),
                )

                # Write updated presentation.xml
                new_prs_xml = self._update_presentation_xml(prs_xml, new_slide_rids)
                dst_zip.writestr(
                    "ppt/presentation.xml",
                    etree.tostring(
                        new_prs_xml,
                        xml_declaration=True,
                        encoding="UTF-8",
                        standalone=True,
                    ),
                )

                # Write updated presentation.xml.rels
                new_prs_rels = self._update_presentation_rels(
                    prs_rels, new_slide_rids, new_slide_targets
                )
                dst_zip.writestr(
                    "ppt/_rels/presentation.xml.rels",
                    etree.tostring(
                        new_prs_rels,
                        xml_declaration=True,
                        encoding="UTF-8",
                        standalone=True,
                    ),
                )

        return tmp_path

    def _update_presentation_xml(self, original_xml, new_rids: list[str]):
        """Replace sldIdLst in presentation.xml with new slide references."""
        xml      = copy.deepcopy(original_xml)
        sldIdLst = xml.find(f".//{{{NS_PML}}}sldIdLst")
        if sldIdLst is None:
            return xml
        for child in list(sldIdLst):
            sldIdLst.remove(child)
        for i, rId in enumerate(new_rids):
            elem = etree.SubElement(sldIdLst, f"{{{NS_PML}}}sldId")
            elem.set("id", str(256 + i))
            elem.set(f"{{{NS_RID}}}id", rId)
        return xml

    def _update_presentation_rels(
        self, original_rels, new_rids: list[str], new_targets: list[str]
    ):
        """Replace slide relationships in presentation.xml.rels."""
        SLIDE_REL_TYPE = (
            "http://schemas.openxmlformats.org/officeDocument/2006/"
            "relationships/slide"
        )
        rels = copy.deepcopy(original_rels)
        for rel in list(rels.findall(f"{{{NS_REL}}}Relationship")):
            if rel.get("Type") == SLIDE_REL_TYPE:
                rels.remove(rel)
        for rId, target in zip(new_rids, new_targets):
            elem = etree.SubElement(rels, f"{{{NS_REL}}}Relationship")
            elem.set("Id",     rId)
            elem.set("Type",   SLIDE_REL_TYPE)
            elem.set("Target", target)
        return rels

    # â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•
    # V1 content-replacement helpers  (unchanged)
    # â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•

    def _apply_content(self, slide, category: str, content: dict):
        """Replace dynamic text fields in the slide (V1 mechanism)."""
        if category not in self.category_map:
            return
        slide_info = self.category_map[category]
        for field in slide_info.get("dynamic_fields", []):
            field_name = field["name"]
            shape_id   = field.get("shape_id")
            if field_name not in content:
                continue
            new_text = str(content[field_name])
            shape    = self._find_shape_by_id(slide, shape_id)
            if shape and shape.has_text_frame:
                self._set_text_preserving_format(shape, new_text)

    def _find_shape_by_id(self, slide, shape_id: int):
        for shape in slide.shapes:
            if shape.shape_id == shape_id:
                return shape
        return None

    def _set_text_preserving_format(self, shape, new_text: str):
        """Replace text preserving original font/size/color/alignment."""
        if not shape.has_text_frame:
            return
        tf    = shape.text_frame
        lines = new_text.split("|")

        ref_run_xml  = None
        ref_para_xml = None
        if tf.paragraphs and tf.paragraphs[0].runs:
            ref_run_xml = copy.deepcopy(tf.paragraphs[0].runs[0]._r)
        if tf.paragraphs:
            ref_para_xml = copy.deepcopy(tf.paragraphs[0]._p)

        txBody = tf._txBody
        paras  = txBody.findall(qn("a:p"))
        for p in paras[1:]:
            txBody.remove(p)

        self._set_paragraph_text(tf.paragraphs[0], lines[0], ref_run_xml)

        for line in lines[1:]:
            if ref_para_xml is not None:
                new_p = copy.deepcopy(ref_para_xml)
                for r in new_p.findall(qn("a:r")):
                    new_p.remove(r)
            else:
                new_p = etree.SubElement(txBody, qn("a:p"))

            if ref_run_xml is not None:
                new_r  = copy.deepcopy(ref_run_xml)
            else:
                new_r  = etree.SubElement(new_p, qn("a:r"))

            t_elem = new_r.find(qn("a:t"))
            if t_elem is None:
                t_elem = etree.SubElement(new_r, qn("a:t"))
            t_elem.text = line
            new_p.append(new_r)
            txBody.append(new_p)

    def _set_paragraph_text(self, para, text: str, ref_run_xml=None):
        """Set text of a single paragraph, preserving run formatting."""
        runs = para.runs
        for run in runs[1:]:
            run._r.getparent().remove(run._r)
        if runs:
            runs[0].text = text
        else:
            p_elem = para._p
            if ref_run_xml is not None:
                new_r  = copy.deepcopy(ref_run_xml)
                t      = new_r.find(qn("a:t"))
                if t is None:
                    t  = etree.SubElement(new_r, qn("a:t"))
                t.text = text
                p_elem.append(new_r)
            else:
                new_r  = etree.SubElement(p_elem, qn("a:r"))
                t      = etree.SubElement(new_r, qn("a:t"))
                t.text = text

    # â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•
    # Utility
    # â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•

    def get_available_categories(self) -> list[str]:
        return [s["category"] for s in self.brand["slides"]]

    def get_category_info(self, category: str) -> dict:
        return self.category_map.get(category, {})


# â”€â”€ Smoke test â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€

if __name__ == "__main__":
    import sys
    sys.stdout.reconfigure(encoding="utf-8")

    PROJECT = r"C:\Users\uilstbadrakh\OneDrive - InvesCore\Desktop\Invescore Slide Generator"
    engine  = InvescoreTemplateEngine(
        os.path.join(PROJECT, "backend", "templates", "InvesCore_Master_Template.pptx"),
        os.path.join(PROJECT, "backend", "brand_guide.json"),
    )

    print("Available categories:", engine.get_available_categories())
    print("Content area bounds :", engine.get_content_area_bounds())

    # â”€â”€ Test V2 with hardcoded content code â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
    HARDCODED_CODE = '''
def build_content(slide, Inches, Pt, Emu, RGBColor):
    """Test: three KPI metric cards."""
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    # Slide title
    tb = slide.shapes.add_textbox(Inches(0.57), Inches(1.35), Inches(8.0), Inches(0.45))
    tf = tb.text_frame
    p  = tf.paragraphs[0]
    r  = p.add_run()
    r.text           = "KEY PERFORMANCE INDICATORS"
    r.font.name      = "Montserrat"
    r.font.size      = Pt(15)
    r.font.bold      = True
    r.font.color.rgb = RGBColor(0x3B, 0x3B, 0x3B)

    # Three metric cards
    metrics = [
        ("\u20ae15.2B", "Revenue",      "+18% YoY"),
        ("23%",         "Market Share", "+2.1pp"),
        ("98.5%",       "Occupancy",    "-0.3pp"),
    ]

    card_w = Inches(2.70)
    card_h = Inches(2.80)
    gap    = Inches(0.20)
    top    = Inches(1.95)
    left0  = Inches(0.57)

    for i, (value, label, delta) in enumerate(metrics):
        left = left0 + i * (card_w + gap)

        # Card background
        bg = slide.shapes.add_shape(1, left, top, card_w, card_h)
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
        bg.line.fill.background()

        # Big number
        nb = slide.shapes.add_textbox(
            left + Inches(0.2), top + Inches(0.3), card_w - Inches(0.4), Inches(0.85)
        )
        tf = nb.text_frame
        p  = tf.paragraphs[0]
        r  = p.add_run()
        r.text           = value
        r.font.name      = "Montserrat"
        r.font.size      = Pt(28)
        r.font.bold      = True
        r.font.color.rgb = RGBColor(0x0C, 0x29, 0x3B)

        # Label
        lb = slide.shapes.add_textbox(
            left + Inches(0.2), top + Inches(1.25), card_w - Inches(0.4), Inches(0.35)
        )
        tf = lb.text_frame
        p  = tf.paragraphs[0]
        r  = p.add_run()
        r.text           = label.upper()
        r.font.name      = "Montserrat"
        r.font.size      = Pt(8)
        r.font.color.rgb = RGBColor(0xA0, 0xAC, 0xBD)

        # Delta
        db = slide.shapes.add_textbox(
            left + Inches(0.2), top + Inches(1.70), card_w - Inches(0.4), Inches(0.30)
        )
        tf = db.text_frame
        p  = tf.paragraphs[0]
        r  = p.add_run()
        r.text           = delta
        r.font.name      = "Montserrat"
        r.font.size      = Pt(9)
        r.font.color.rgb = (
            RGBColor(0xC8, 0x10, 0x2E)
            if delta.startswith("+")
            else RGBColor(0x66, 0x66, 0x66)
        )
'''

    slide_plan = {
        "presentation_title": "PHASE 1 BRAND-FRAME TEST",
        "sections": [
            {
                "name": "EXECUTIVE SUMMARY",
                "slides": [
                    {
                        "slide_type": "content",
                        "title": "Key Performance Indicators",
                        "description": "Three KPI metric cards showing revenue, market share, occupancy.",
                        "content_spec": {},
                    }
                ],
            },
            {
                "name": "MARKET ANALYSIS",
                "slides": [
                    {
                        "slide_type": "content",
                        "title": "Second Slide (fallback test â€” no code)",
                        "description": "Should show brand frame + fallback title only.",
                        "content_spec": {},
                    }
                ],
            },
        ],
    }

    # Slide 0 gets the metric-card code; slide 1 intentionally left blank.
    content_code_map = {0: HARDCODED_CODE}

    out, _warnings = engine.create_presentation_v2(slide_plan, content_code_map)
    dest = os.path.join(PROJECT, "backend", "phase1_test_output.pptx")
    shutil.move(out, dest)
    print(f"\nPhase 1 test output: {dest}")
    print("Open the file and verify:")
    print("  Slide 1 â€” Opening: 'PHASE 1 BRAND-FRAME TEST'")
    print("  Slide 2 â€” Agenda:  two sections with page ranges")
    print("  Slide 3 â€” Content: brand frame + 3 KPI metric cards")
    print("  Slide 4 â€” Content: brand frame + fallback title (no AI code)")
    print("  Slide 5 â€” Closing: unchanged clone")

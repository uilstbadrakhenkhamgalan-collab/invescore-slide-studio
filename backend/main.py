"""
InvesCore Slide Studio — FastAPI Backend v2
"""
import asyncio
import hashlib
import json
import logging
import os
import re
import secrets
import shutil
import time
import uuid
from contextlib import asynccontextmanager
from pathlib import Path
from typing import Any, AsyncGenerator

import anthropic
from fastapi import FastAPI, Header, HTTPException, Request
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, JSONResponse, StreamingResponse
from pydantic import BaseModel, Field, field_validator
from starlette.middleware.base import BaseHTTPMiddleware

from template_engine import InvescoreTemplateEngine

# ── Logging ───────────────────────────────────────────────────────────────────
# Two goals:
#   1. Always include a request_id field for correlation across endpoints.
#   2. Never log a raw API key, even if it slips into an error message.

_API_KEY_RE_LOG = re.compile(r"sk-ant-[A-Za-z0-9_\-]{6,}")


class RedactingFormatter(logging.Formatter):
    """Strip Anthropic API keys from any formatted log line."""

    def format(self, record: logging.LogRecord) -> str:
        msg = super().format(record)
        return _API_KEY_RE_LOG.sub("sk-ant-***REDACTED***", msg)


def _configure_logging() -> None:
    root = logging.getLogger()
    if getattr(root, "_invescore_configured", False):
        return
    handler = logging.StreamHandler()
    handler.setFormatter(RedactingFormatter(
        "%(asctime)s %(levelname)s %(name)s %(message)s",
        "%Y-%m-%dT%H:%M:%S",
    ))
    root.handlers[:] = [handler]
    root.setLevel(os.environ.get("LOG_LEVEL", "INFO").upper())
    # Quiet uvicorn's noisy access log unless explicitly debugging.
    logging.getLogger("uvicorn.access").setLevel(
        os.environ.get("UVICORN_ACCESS_LEVEL", "WARNING").upper()
    )
    root._invescore_configured = True  # type: ignore[attr-defined]


_configure_logging()
logger = logging.getLogger("invescore.api")

# ── Paths ─────────────────────────────────────────────────────────────────────
BASE_DIR        = Path(__file__).parent
TEMPLATE_PATH   = BASE_DIR / "templates" / "InvesCore_Master_Template.pptx"
BRAND_GUIDE_PATH = BASE_DIR / "brand_guide.json"
TMP_DIR         = BASE_DIR / "tmp"
ARTIFACT_TTL_SECONDS = int(os.environ.get("ARTIFACT_TTL_SECONDS", "21600"))

# ── Limits (defense-in-depth against DoS-via-cost on user keys) ───────────────
MAX_REQUEST_BYTES   = int(os.environ.get("MAX_REQUEST_BYTES", str(256 * 1024)))  # 256 KB
MAX_SECTIONS        = 8           # template nav-bar has 8 slots
MAX_CONTENT_SLIDES  = 25          # interpreter prompt allows 5–20; this is a hard cap
MAX_INTAKE_TURNS    = 30
MAX_INTAKE_MSG_LEN  = 8_000
MAX_PROMPT_LEN      = 16_000

# Builder concurrency: how many slides to call in parallel. Anthropic Tier 1 is
# 5 RPM input; Tier 2 is 1k RPM. Default to 6 — fast for most users, retry-able
# if rate-limited. Override with BUILDER_CONCURRENCY env var.
BUILDER_CONCURRENCY = int(os.environ.get("BUILDER_CONCURRENCY", "6"))
BUILDER_RETRIES     = int(os.environ.get("BUILDER_RETRIES", "3"))
SSE_HEARTBEAT_SEC   = int(os.environ.get("SSE_HEARTBEAT_SEC", "15"))

# ── App ───────────────────────────────────────────────────────────────────────
@asynccontextmanager
async def _lifespan(app: FastAPI):
    # Startup: validate environment + launch periodic cleanup.
    if not TEMPLATE_PATH.exists():
        raise RuntimeError(f"Master template missing: {TEMPLATE_PATH}")
    if not BRAND_GUIDE_PATH.exists():
        raise RuntimeError(f"Brand guide missing: {BRAND_GUIDE_PATH}")
    try:
        from pptx import Presentation as _P
        _P(str(TEMPLATE_PATH))
    except Exception as exc:
        raise RuntimeError(f"Master template is not a valid pptx: {exc}")
    cleanup_task = asyncio.create_task(_periodic_cleanup())
    logger.info("startup ok — template=%s cors=%s", TEMPLATE_PATH.name, ALLOWED_ORIGINS)
    try:
        yield
    finally:
        cleanup_task.cancel()
        try:
            await cleanup_task
        except (asyncio.CancelledError, Exception):
            pass


app = FastAPI(title="InvesCore Slide Studio API", version="2.1.0", lifespan=_lifespan)


class ContentLengthLimitMiddleware(BaseHTTPMiddleware):
    """Reject oversized request bodies before any body is read."""

    def __init__(self, app, max_bytes: int):
        super().__init__(app)
        self.max_bytes = max_bytes

    async def dispatch(self, request: Request, call_next):
        cl = request.headers.get("content-length")
        if cl is not None:
            try:
                if int(cl) > self.max_bytes:
                    return JSONResponse(
                        status_code=413,
                        content={"detail": "Request body too large"},
                    )
            except ValueError:
                return JSONResponse(
                    status_code=400,
                    content={"detail": "Invalid Content-Length"},
                )
        return await call_next(request)


class RequestIdMiddleware(BaseHTTPMiddleware):
    """Tag every request with a 12-char id; surface it in X-Request-ID header."""

    async def dispatch(self, request: Request, call_next):
        rid = request.headers.get("x-request-id") or uuid.uuid4().hex[:12]
        request.state.request_id = rid
        start = time.monotonic()
        response = await call_next(request)
        elapsed_ms = int((time.monotonic() - start) * 1000)
        response.headers["X-Request-ID"] = rid
        # Skip noisy paths from access log
        if request.url.path != "/api/health":
            logger.info(
                "req rid=%s method=%s path=%s status=%s ms=%d",
                rid, request.method, request.url.path, response.status_code, elapsed_ms,
            )
        return response


app.add_middleware(RequestIdMiddleware)
app.add_middleware(ContentLengthLimitMiddleware, max_bytes=MAX_REQUEST_BYTES)

# CORS allow-list is environment-driven. Defaults match the deployed Vercel
# domain + local dev. Wildcards rejected — a user-supplied API key flowing
# through this backend means an open CORS = open AI-cost laundering proxy.
_default_origins = ",".join([
    "https://invescore-slide-studio.vercel.app",
    "http://localhost:3000",
    "http://127.0.0.1:3000",
])
_origins_raw = os.environ.get("CORS_ALLOW_ORIGINS", _default_origins)
ALLOWED_ORIGINS = [o.strip() for o in _origins_raw.split(",") if o.strip() and o.strip() != "*"]

app.add_middleware(
    CORSMiddleware,
    allow_origins=ALLOWED_ORIGINS,
    allow_credentials=False,
    allow_methods=["GET", "POST", "OPTIONS"],
    allow_headers=["Content-Type", "X-Download-Token"],
)

# ── Pydantic models ───────────────────────────────────────────────────────────
_ANTHROPIC_KEY_RE = re.compile(r"^sk-ant-[A-Za-z0-9_\-]{20,}$")


def _validate_api_key(value: str) -> str:
    if not isinstance(value, str) or not _ANTHROPIC_KEY_RE.match(value.strip()):
        raise ValueError("Invalid Anthropic API key format")
    return value.strip()


class InterpretRequest(BaseModel):
    api_key: str
    prompt: str = Field(..., min_length=1, max_length=MAX_PROMPT_LEN)

    @field_validator("api_key")
    @classmethod
    def _check_key(cls, v: str) -> str:
        return _validate_api_key(v)


class GenerateRequest(BaseModel):
    api_key: str
    slide_plan: list[dict]   # V1: flat list of {template, content}

    @field_validator("api_key")
    @classmethod
    def _check_key(cls, v: str) -> str:
        return _validate_api_key(v)


class GenerateV2Request(BaseModel):
    api_key: str
    slide_plan: dict         # V2: {presentation_title, sections:[...]}

    @field_validator("api_key")
    @classmethod
    def _check_key(cls, v: str) -> str:
        return _validate_api_key(v)


class IntakeRequest(BaseModel):
    api_key: str
    messages: list[dict]     # [{role, content}, ...]

    @field_validator("api_key")
    @classmethod
    def _check_key(cls, v: str) -> str:
        return _validate_api_key(v)


def _safe_error_response(exc: Exception, request_id: str, default: str) -> dict:
    """
    Return a sanitized error payload. The full exception is logged server-side
    with the request_id so support can correlate; the client gets a generic
    message + the id (no stack traces, no internal paths).
    """
    logger.exception("[%s] %s: %s", request_id, default, exc)
    return {"detail": f"{default} (ref: {request_id})"}


def _validate_v2_slide_plan(plan: dict) -> tuple[bool, str]:
    """Hard caps on the V2 slide plan to bound AI-cost and template overflow."""
    if not isinstance(plan, dict):
        return False, "slide_plan must be an object"
    sections = plan.get("sections")
    if not isinstance(sections, list) or not sections:
        return False, "slide_plan.sections must be a non-empty list"
    if len(sections) > MAX_SECTIONS:
        return False, f"max {MAX_SECTIONS} sections allowed (got {len(sections)})"
    total_content = 0
    for sec in sections:
        if not isinstance(sec, dict) or "name" not in sec:
            return False, "each section must have a `name`"
        if not isinstance(sec.get("name"), str) or not sec["name"].strip():
            return False, "each section name must be a non-empty string"
        slides = sec.get("slides", [])
        if not isinstance(slides, list):
            return False, "section.slides must be a list"
        for s in slides:
            if not isinstance(s, dict):
                return False, "each slide entry must be an object"
            if s.get("slide_type", "content") == "content":
                total_content += 1
    if total_content == 0:
        return False, "slide_plan must contain at least one content slide"
    if total_content > MAX_CONTENT_SLIDES:
        return False, (
            f"max {MAX_CONTENT_SLIDES} content slides allowed "
            f"(got {total_content})"
        )
    title = plan.get("presentation_title", "")
    if not isinstance(title, str) or len(title) > 300:
        return False, "presentation_title must be a string <= 300 chars"
    return True, ""

# ── Load brand guide ──────────────────────────────────────────────────────────
with open(BRAND_GUIDE_PATH, encoding="utf-8") as f:
    _brand = json.load(f)

# ── Content-area bounds (passed to Builder Agent) ─────────────────────────────
_engine_tmp = InvescoreTemplateEngine(str(TEMPLATE_PATH), str(BRAND_GUIDE_PATH))
_CONTENT_BOUNDS = _engine_tmp.get_content_area_bounds()

# ═════════════════════════════════════════════════════════════════════════════
# SYSTEM PROMPTS
# ═════════════════════════════════════════════════════════════════════════════

INTAKE_SYSTEM_PROMPT = """You are the InvesCore Slide Studio intake assistant. Your job is to interview the user to gather everything needed to create a perfect branded presentation.

Ask questions ONE AT A TIME. Keep them short and conversational. After each answer, acknowledge briefly and ask the next question. Adapt your follow-up questions based on their answers.

If the user's first message is exactly "START", greet them briefly and ask the first question.

REQUIRED INFORMATION TO GATHER (in roughly this order):

1. TOPIC & PURPOSE
   "What is this presentation about?"
   Then follow up: "Who is the audience? (e.g., internal team, investors, clients, board)"

2. LANGUAGE
   "Should the slides be in English, Mongolian, or a mix of both?"

3. SCOPE & LENGTH
   "How many slides do you have in mind? Or should I decide based on the content?"
   (If they say "you decide", suggest a number based on the topic complexity)

4. SECTIONS / STRUCTURE
   "What main sections or topics should the presentation cover?"
   (If they're unsure, suggest logical sections based on the topic and ask them to confirm or adjust)

5. KEY DATA & NUMBERS
   "Do you have any specific numbers, statistics, or data points you want included?"
   "Any specific dates, deadlines, or time periods to reference?"

6. TONE & STYLE
   "What tone should this have? (e.g., formal/executive, informational, persuasive, casual update)"

7. SPECIAL REQUIREMENTS
   "Anything else I should know? Any slides that need a specific layout like a table, chart, comparison, or org chart?"

RULES:
- Never ask more than 2 questions in a single message
- If the user gives a detailed answer that covers multiple questions, skip the ones already answered
- If the user says "that's it" or "nothing else" or seems done, stop asking and confirm
- Keep your tone professional but warm, like a helpful colleague
- After gathering everything, send a final confirmation message summarizing the plan
- Total conversation should be 5-8 exchanges, not more
- If the user's very first message is already detailed (200+ characters), skip most questions and just confirm

WHEN YOU HAVE ENOUGH INFORMATION, end your response with exactly this format (after your conversational sign-off):

---INTAKE COMPLETE---
{
  "topic": "...",
  "audience": "...",
  "language": "...",
  "slide_count": "...",
  "sections": ["...", "..."],
  "key_data": ["...", "..."],
  "tone": "...",
  "special_requests": "...",
  "full_brief": "A comprehensive paragraph combining all gathered info into a detailed prompt for the Interpreter Agent"
}
---END---

The full_brief field is the most important — it should read like a detailed, specific presentation request that incorporates every piece of information the user provided."""


INTERPRETER_SYSTEM_PROMPT = """You are the InvesCore Slide Studio Interpreter. Convert a presentation brief into a structured slide plan for a professional Mongolian real estate and investment company.

IMPORTANT — OUTPUT FORMAT:
Return ONLY valid JSON, no markdown fences, no explanation. The schema:

{
  "presentation_title": "TITLE IN CAPS",
  "sections": [
    {
      "name": "SECTION NAME IN CAPS",
      "slides": [
        {
          "slide_type": "content",
          "title": "SLIDE TITLE",
          "description": "Detailed description of what this slide shows, what layout works best, what visual elements to include. Be specific. This directly guides the Builder Agent that will write python-pptx code.",
          "content_spec": {
            "layout": "metric_callout | two_column | table_focus | chart_focus | timeline | comparison | bullets | mixed | freeform",
            "elements": [
              { "type": "title", "text": "THE SLIDE TITLE" },
              { "type": "subtitle", "text": "Optional subtitle or section label" },
              { "type": "bullets", "items": ["Point 1", "Point 2", "Point 3"] },
              {
                "type": "table",
                "headers": ["Col A", "Col B", "Col C"],
                "rows": [["R1A","R1B","R1C"], ["R2A","R2B","R2C"]]
              },
              {
                "type": "chart",
                "chart_type": "bar | column | line | pie | stacked_bar",
                "title": "Chart Title",
                "categories": ["Cat1","Cat2","Cat3"],
                "series": [{"name": "Series1", "values": [10,20,30]}]
              },
              { "type": "metric", "value": "₮15.2B", "label": "Revenue", "delta": "+18% YoY" },
              { "type": "text_block", "heading": "Key Insight", "body": "..." },
              { "type": "callout", "text": "Important highlighted statement" },
              {
                "type": "comparison",
                "left": {"title": "Option A", "points": ["..."]},
                "right": {"title": "Option B", "points": ["..."]}
              },
              {
                "type": "timeline",
                "steps": [{"label": "Q1 2025", "description": "Phase 1 launch"}]
              },
              {
                "type": "diagram_description",
                "description": "Describe a visual: e.g. '3-step flow: Acquisition → Development → Exit, each in a dark navy box with connecting arrows'"
              }
            ]
          }
        },
        {
          "slide_type": "section_divider",
          "title": "SECTION TITLE",
          "description": "One-sentence overview of this section"
        }
      ]
    }
  ]
}

RULES:
1. Opening, agenda, and closing slides are automatic — do NOT include them
2. Define sections (max 8) — the agenda is auto-generated from them
3. Each section: 1–5 slides
4. Total: 5–20 content slides
5. Write ALL actual text, numbers, and data — be specific; use plausible InvesCore data
6. Choose the right visual for each slide's purpose:
   - Financial data → table or chart
   - KPIs / key numbers → metric_callout (big numbers, card layout)
   - Strategy / narrative → bullets or text_block
   - Process / phases → timeline or diagram_description
   - Side-by-side alternatives → comparison
   - Mixed data + commentary → mixed or two_column
7. Vary layouts — do not repeat the same layout on consecutive slides
8. The description field is critical — the Builder Agent reads it to decide how to lay out shapes
9. Support Mongolian (Cyrillic) content when language is Mongolian
10. Think like an investment analyst preparing a board-level presentation"""


BUILDER_SYSTEM_PROMPT = f"""You are the InvesCore Slide Studio Builder. You generate python-pptx code that creates the CONTENT AREA of a single PowerPoint slide.

SLIDE DIMENSIONS: 10 inches wide × 5.625 inches tall.

YOUR CONTENT AREA (stay strictly within these bounds):
  Top:    {_CONTENT_BOUNDS['top']} inches from slide top
  Bottom: {_CONTENT_BOUNDS['bottom']} inches from slide top
  Left:   {_CONTENT_BOUNDS['left']} inches from slide left
  Right:  {_CONTENT_BOUNDS['right']} inches from slide left
  → Working area ≈ {_CONTENT_BOUNDS['right'] - _CONTENT_BOUNDS['left']:.2f}" wide × {_CONTENT_BOUNDS['bottom'] - _CONTENT_BOUNDS['top']:.2f}" tall

The brand frame (header bar, logo, nav bar, page number) is already on the slide — do NOT add these.

BRAND COLORS (use ONLY these):
  PRIMARY_DARK = RGBColor(0x3B, 0x3B, 0x3B)   # main text
  WHITE        = RGBColor(0xFF, 0xFF, 0xFF)    # text on dark bg
  ACCENT_RED   = RGBColor(0xC8, 0x10, 0x2E)   # emphasis only, use sparingly
  LIGHT_GRAY   = RGBColor(0xA0, 0xAC, 0xBD)   # secondary / labels
  DARK_NAVY    = RGBColor(0x0C, 0x29, 0x3B)   # dark card backgrounds
  MEDIUM_GRAY  = RGBColor(0x66, 0x66, 0x66)   # supporting text
  LIGHT_BG     = RGBColor(0xF5, 0xF5, 0xF5)   # subtle backgrounds
  BORDER_GRAY  = RGBColor(0xE0, 0xE0, 0xE0)   # borders / dividers

FONTS (Montserrat only):
  Section label  :  9 pt, Regular,  LIGHT_GRAY   (e.g. "EXECUTIVE SUMMARY")
  Slide title    : 13 pt, SemiBold, PRIMARY_DARK
  Subtitle       :  9 pt, Regular,  MEDIUM_GRAY
  Body / bullets :  7 pt, Regular,  PRIMARY_DARK
  Caption        :  6 pt, Regular,  MEDIUM_GRAY
  Big metric     : 22–28 pt, Bold,  DARK_NAVY or ACCENT_RED

QUALITY STANDARDS:
- Think like a McKinsey/Goldman designer — clear hierarchy, intentional whitespace
- Tables: dark navy header row (white text), alternating LIGHT_BG / WHITE rows
- Charts: brand colors for series, no gridlines unless essential, minimal axes
- Metrics: large font (22–28 pt), card background (LIGHT_BG), label in LIGHT_GRAY below
- Layouts must feel balanced — do not cram everything in
- No drop shadows, no 3D, no WordArt

OUTPUT FORMAT — return ONLY this Python function, no explanation, no markdown fences:

def build_content(slide, Inches, Pt, Emu, RGBColor):
    \"\"\"Build content area.\"\"\"
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    # ... your code ...

IMPORTANT:
- All imports must be INSIDE the function
- Use only: pptx, math, datetime, random, itertools, collections, decimal, statistics
- Do NOT use: os, sys, subprocess, open(), exec(), eval(), requests, socket, pathlib
- The function will be executed with exec() — it must be self-contained
- Shape type constant for rectangles: use 1 (MSO_SHAPE_TYPE.RECTANGLE equivalent)
- For charts, import from pptx.chart.data import ChartData and from pptx.enum.chart import XL_CHART_TYPE
- Always add a slide title text box as the very first element

EXAMPLE — metric callout slide:

def build_content(slide, Inches, Pt, Emu, RGBColor):
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    # Title
    tb = slide.shapes.add_textbox(Inches(0.57), Inches(1.35), Inches(7.5), Inches(0.40))
    tf = tb.text_frame
    p  = tf.paragraphs[0]
    r  = p.add_run()
    r.text = "KEY PERFORMANCE INDICATORS"
    r.font.name = "Montserrat"
    r.font.size = Pt(13)
    r.font.bold = True
    r.font.color.rgb = RGBColor(0x3B, 0x3B, 0x3B)

    metrics = [
        ("\\u20ae15.2B", "Revenue",      "+18% YoY"),
        ("23%",          "Market Share", "+2.1pp"),
        ("98.5%",        "Occupancy",    "-0.3pp"),
    ]
    card_w = Inches(2.50); card_h = Inches(2.60)
    gap    = Inches(0.18); top    = Inches(1.88); left0 = Inches(0.57)

    for i, (value, label, delta) in enumerate(metrics):
        left = left0 + i * (card_w + gap)
        bg = slide.shapes.add_shape(1, left, top, card_w, card_h)
        bg.fill.solid(); bg.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
        bg.line.fill.background()

        nb = slide.shapes.add_textbox(left+Inches(0.18), top+Inches(0.28), card_w-Inches(0.36), Inches(0.72))
        tf = nb.text_frame; p = tf.paragraphs[0]; r = p.add_run()
        r.text = value; r.font.name = "Montserrat"; r.font.size = Pt(26)
        r.font.bold = True; r.font.color.rgb = RGBColor(0x0C, 0x29, 0x3B)

        lb = slide.shapes.add_textbox(left+Inches(0.18), top+Inches(1.12), card_w-Inches(0.36), Inches(0.28))
        tf = lb.text_frame; p = tf.paragraphs[0]; r = p.add_run()
        r.text = label.upper(); r.font.name = "Montserrat"; r.font.size = Pt(7)
        r.font.color.rgb = RGBColor(0xA0, 0xAC, 0xBD)

        db = slide.shapes.add_textbox(left+Inches(0.18), top+Inches(1.52), card_w-Inches(0.36), Inches(0.25))
        tf = db.text_frame; p = tf.paragraphs[0]; r = p.add_run()
        r.text = delta; r.font.name = "Montserrat"; r.font.size = Pt(8)
        r.font.color.rgb = RGBColor(0xC8,0x10,0x2E) if delta.startswith("+") else RGBColor(0x66,0x66,0x66)"""


# ═════════════════════════════════════════════════════════════════════════════
# Helper — call Builder Agent (Opus)
# ═════════════════════════════════════════════════════════════════════════════

def _call_builder_agent(
    api_key: str,
    slide_spec: dict,
    presentation_title: str,
    section_name: str,
    all_sections: list[str],
) -> str:
    """
    Calls claude-opus-4-6 to generate a build_content() python-pptx function
    for a single content slide. Returns the raw code string.
    """
    client = anthropic.Anthropic(api_key=api_key)

    user_msg = f"""Generate build_content() for this slide:

PRESENTATION: {presentation_title}
SECTION: {section_name}
ALL SECTIONS: {', '.join(all_sections)}

SLIDE SPEC:
{json.dumps(slide_spec, indent=2, ensure_ascii=False)}

Return ONLY the Python function — no explanation, no markdown fences."""

    # Prompt caching: the system prompt is identical across every slide in a
    # generation, so cache_control: ephemeral makes calls 2..N read 90% cheaper
    # and ~150ms faster TTFT.
    response = client.messages.create(
        model="claude-opus-4-6",
        max_tokens=16000,
        temperature=0.2,
        system=[
            {
                "type": "text",
                "text": BUILDER_SYSTEM_PROMPT,
                "cache_control": {"type": "ephemeral"},
            }
        ],
        messages=[{"role": "user", "content": user_msg}],
    )

    if not response.content:
        raise RuntimeError("Builder Agent returned an empty response")
    code = response.content[0].text
    return _extract_python_code(code)


_FENCED_RE = re.compile(r"```(?:python)?\s*\n(.*?)```", re.DOTALL)
_DEF_RE    = re.compile(r"(?m)^\s*def\s+build_content\s*\(")


def _extract_python_code(text: str) -> str:
    """
    Pull a `def build_content(...)` block out of a noisy LLM response.

    Handles:
      - clean function returned as-is
      - ```python ... ``` fenced blocks (anywhere in the text)
      - prose preambles ('Here is the code:') before the function
    """
    text = text.strip()
    # Try a fenced block first — it usually contains exactly the function.
    fenced_matches = _FENCED_RE.findall(text)
    for fenced in fenced_matches:
        if _DEF_RE.search(fenced):
            return fenced.strip()
    # No fences — drop any preamble before the first `def build_content`.
    m = _DEF_RE.search(text)
    if m:
        return text[m.start():].strip()
    return text.strip()


# ═════════════════════════════════════════════════════════════════════════════
# SSE helper
# ═════════════════════════════════════════════════════════════════════════════

def _sse(event: str, data: Any = None) -> str:
    """Format a Server-Sent Event string."""
    line = f"event: {event}\n"
    line += f"data: {json.dumps(data) if data is not None else ''}\n\n"
    return line


def _run_intake_request(api_key: str, messages: list[dict]) -> str:
    client = anthropic.Anthropic(api_key=api_key)
    message = client.messages.create(
        model="claude-sonnet-4-6",
        max_tokens=2000,
        temperature=0.4,
        system=[
            {
                "type": "text",
                "text": INTAKE_SYSTEM_PROMPT,
                "cache_control": {"type": "ephemeral"},
            }
        ],
        messages=messages,
    )
    if not message.content:
        raise RuntimeError("Intake Agent returned an empty response")
    return message.content[0].text


def _run_interpreter_request(api_key: str, prompt: str):
    client = anthropic.Anthropic(api_key=api_key)
    with client.messages.stream(
        model="claude-sonnet-4-6",
        max_tokens=32000,
        temperature=0.3,
        system=[
            {
                "type": "text",
                "text": INTERPRETER_SYSTEM_PROMPT,
                "cache_control": {"type": "ephemeral"},
            }
        ],
        messages=[{"role": "user", "content": prompt}],
    ) as stream:
        raw = stream.get_final_text()
        usage = stream.get_final_message().usage
    return raw, usage


def _sanitize_download_filename(raw_title: str) -> str:
    """
    Produce a safe Content-Disposition filename from an untrusted title.
    Strips control chars, path separators, dot-runs (.., ...), and reserved
    Windows characters. Length-capped to keep most clients happy.
    """
    if not isinstance(raw_title, str):
        raw_title = ""
    cleaned = re.sub(r"[\x00-\x1f\x7f]+", "", raw_title).strip()
    # Replace path / drive separators and Windows-reserved characters.
    for ch in ("/", "\\", ":", "*", "?", '"', "<", ">", "|"):
        cleaned = cleaned.replace(ch, "-")
    # Collapse any run of dots so "../.." can't survive in the filename.
    cleaned = re.sub(r"\.{2,}", ".", cleaned)
    cleaned = re.sub(r"\s+", " ", cleaned).strip(" .-")
    if not cleaned:
        cleaned = "Presentation"
    return f"InvesCore_{cleaned[:80]}.pptx"


def _artifact_file_path(artifact_id: str) -> Path:
    return TMP_DIR / f"{artifact_id}.pptx"


def _artifact_meta_path(artifact_id: str) -> Path:
    return TMP_DIR / f"{artifact_id}.json"


def _hash_download_token(token: str) -> str:
    return hashlib.sha256(token.encode("utf-8")).hexdigest()


def _cleanup_expired_artifacts() -> int:
    """
    Sweep expired artifacts. Best-effort: never raises. Returns count removed.

    Uses an in-memory grace window so we don't accidentally delete an artifact
    being written *right now* (race between metadata write and the file move).
    """
    TMP_DIR.mkdir(parents=True, exist_ok=True)
    current_ts = time.time()
    removed = 0

    for meta_path in TMP_DIR.glob("*.json"):
        try:
            meta = json.loads(meta_path.read_text(encoding="utf-8"))
            expires_at = float(meta.get("expires_at", 0))
            created_at = float(meta.get("created_at", 0))
        except Exception:
            # Unreadable metadata — could be mid-write. Skip; next sweep tries again.
            continue

        # Don't reap something less than 60s old: the file may be in flight.
        if created_at > 0 and (current_ts - created_at) < 60:
            continue

        artifact_id = meta_path.stem
        file_path = _artifact_file_path(artifact_id)
        if expires_at > current_ts and file_path.exists():
            continue

        try:
            file_path.unlink(missing_ok=True)
            meta_path.unlink(missing_ok=True)
            removed += 1
        except OSError:
            pass  # best-effort
    return removed


async def _periodic_cleanup() -> None:
    while True:
        try:
            removed = await asyncio.to_thread(_cleanup_expired_artifacts)
            if removed:
                logger.info("artifact cleanup removed %d expired files", removed)
        except Exception as e:
            logger.warning("artifact cleanup failed: %s", e)
        await asyncio.sleep(300)  # 5 min


def _store_download_artifact(source_path: str, presentation_title: str) -> dict[str, str]:
    TMP_DIR.mkdir(parents=True, exist_ok=True)
    _cleanup_expired_artifacts()

    artifact_id = uuid.uuid4().hex
    download_token = secrets.token_urlsafe(32)
    filename = _sanitize_download_filename(presentation_title)
    artifact_path = _artifact_file_path(artifact_id)
    meta_path = _artifact_meta_path(artifact_id)

    shutil.move(source_path, artifact_path)
    metadata = {
        "filename": filename,
        "token_sha256": _hash_download_token(download_token),
        "created_at": time.time(),
        "expires_at": time.time() + ARTIFACT_TTL_SECONDS,
    }
    meta_path.write_text(json.dumps(metadata), encoding="utf-8")

    return {
        "artifact_id": artifact_id,
        "download_token": download_token,
        "filename": filename,
    }


# ═════════════════════════════════════════════════════════════════════════════
# ENDPOINTS
# ═════════════════════════════════════════════════════════════════════════════

@app.get("/api/health")
async def health():
    return {
        "status": "ok",
        "template": str(TEMPLATE_PATH.name),
        "version": "2.1.0",
        "limits": {
            "max_sections": MAX_SECTIONS,
            "max_content_slides": MAX_CONTENT_SLIDES,
            "max_request_bytes": MAX_REQUEST_BYTES,
            "builder_concurrency": BUILDER_CONCURRENCY,
        },
    }


# ── Intake ────────────────────────────────────────────────────────────────────

def _sanitize_intake_messages(messages: list[dict]) -> list[dict]:
    """Strip / truncate messages so we never forward unbounded payloads to Anthropic."""
    if len(messages) > MAX_INTAKE_TURNS:
        messages = messages[-MAX_INTAKE_TURNS:]
    cleaned: list[dict] = []
    for m in messages:
        if not isinstance(m, dict):
            continue
        role = m.get("role")
        content = m.get("content")
        if role not in ("user", "assistant"):
            continue
        if not isinstance(content, str):
            continue
        if len(content) > MAX_INTAKE_MSG_LEN:
            content = content[:MAX_INTAKE_MSG_LEN]
        cleaned.append({"role": role, "content": content})
    return cleaned


@app.post("/api/intake")
async def intake_conversation(req: IntakeRequest):
    """Conversational intake agent."""
    request_id = uuid.uuid4().hex[:12]
    if not req.messages:
        raise HTTPException(400, "messages cannot be empty")
    cleaned = _sanitize_intake_messages(req.messages)
    if not cleaned:
        raise HTTPException(400, "messages contained no valid entries")
    try:
        content = await asyncio.to_thread(
            _run_intake_request,
            req.api_key,
            cleaned,
        )
        return {"content": content}
    except anthropic.AuthenticationError:
        raise HTTPException(401, "Invalid API key")
    except anthropic.RateLimitError:
        raise HTTPException(429, "API rate limit exceeded")
    except Exception as e:
        payload = _safe_error_response(e, request_id, "Intake failed")
        raise HTTPException(500, payload["detail"])


# ── Interpret (upgraded prompt, same endpoint) ─────────────────────────────────

@app.post("/api/interpret")
async def interpret(req: InterpretRequest):
    """
    Interpreter Agent — converts brief into structured slide plan (v2 schema).
    """
    request_id = uuid.uuid4().hex[:12]
    try:
        def _is_overloaded(exc: Exception) -> bool:
            if isinstance(exc, anthropic.APIStatusError):
                return getattr(exc, "status_code", None) == 529
            return "overloaded" in str(exc).lower()

        # Retry up to 3 times on overloaded errors with exponential backoff
        max_retries = 3
        last_exc: Exception | None = None
        for attempt in range(max_retries):
            try:
                raw, usage = await asyncio.to_thread(
                    _run_interpreter_request,
                    req.api_key,
                    req.prompt,
                )
                last_exc = None
                break  # success
            except Exception as e:
                if _is_overloaded(e) and attempt < max_retries - 1:
                    last_exc = e
                    await asyncio.sleep(5 * (attempt + 1))
                    continue
                raise
        if last_exc is not None:
            raise HTTPException(503, "Anthropic API is overloaded — please try again in a moment")

        raw = raw.strip()
        # Strip markdown fences
        raw = re.sub(r"^```json\s*", "", raw)
        raw = re.sub(r"^```\s*",     "", raw)
        raw = re.sub(r"\s*```$",     "", raw)
        raw = raw.strip()

        plan = json.loads(raw)

        # Count total content slides for cost estimate
        total_content = sum(
            len(sec.get("slides", []))
            for sec in plan.get("sections", [])
        )

        return {
            "presentation_title": plan.get("presentation_title", "Presentation"),
            "sections":           plan.get("sections", []),
            "token_usage": {
                "input_tokens":       usage.input_tokens,
                "output_tokens":      usage.output_tokens,
                "estimated_cost_usd": round(
                    (usage.input_tokens * 3 +
                     usage.output_tokens * 15) / 1_000_000, 4
                ),
            },
            "total_content_slides": total_content,
            "estimated_builder_cost_usd": round(
                total_content * (2000 * 15 + 3000 * 75) / 1_000_000, 2
            ),
        }

    except json.JSONDecodeError as e:
        logger.warning("[%s] Interpreter returned invalid JSON: %s", request_id, e)
        raise HTTPException(502, f"Interpreter returned invalid JSON (ref: {request_id})")
    except anthropic.AuthenticationError:
        raise HTTPException(401, "Invalid API key")
    except anthropic.RateLimitError:
        raise HTTPException(429, "API rate limit exceeded")
    except Exception as e:
        if _is_overloaded(e):
            raise HTTPException(503, "Anthropic API is overloaded — please try again in a moment")
        payload = _safe_error_response(e, request_id, "Interpretation failed")
        raise HTTPException(500, payload["detail"])


# ── V1 Generate (unchanged — kept for backwards compat) ────────────────────────

@app.post("/api/generate")
async def generate(req: GenerateRequest):
    """V1: Generate .pptx from flat slide_plan (clone + text swap)."""
    request_id = uuid.uuid4().hex[:12]
    if not req.slide_plan:
        raise HTTPException(400, "slide_plan cannot be empty")
    if len(req.slide_plan) > MAX_CONTENT_SLIDES + 4:  # +structural slides
        raise HTTPException(400, f"slide_plan too large (max {MAX_CONTENT_SLIDES + 4} entries)")
    try:
        engine = InvescoreTemplateEngine(str(TEMPLATE_PATH), str(BRAND_GUIDE_PATH))
        output_path = engine.create_presentation(req.slide_plan)

        title = "presentation"
        if req.slide_plan and req.slide_plan[0].get("content", {}).get("presentation_title"):
            raw   = req.slide_plan[0]["content"]["presentation_title"]
            title = raw[:40].replace("/", "-").replace("\\", "-").replace(" ", "_")

        filename = f"InvesCore_{title}.pptx"
        return FileResponse(
            path        = output_path,
            media_type  = "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            filename    = filename,
            background  = None,
        )
    except ValueError as e:
        raise HTTPException(400, str(e))
    except Exception as e:
        payload = _safe_error_response(e, request_id, "Generation failed")
        raise HTTPException(500, payload["detail"])


# ── V2 Generate — SSE streaming, Builder Agent per slide ──────────────────────

async def _builder_with_retry(
    api_key: str,
    slide_spec: dict,
    presentation_title: str,
    section_name: str,
    all_sections: list[str],
    request_id: str,
    cidx: int,
) -> tuple[int, str | None, str | None]:
    """
    Run a single Builder Agent call with bounded retry on overload / rate-limit.
    Returns (cidx, code_or_None, error_str_or_None).
    """
    delay = 2.0
    last_exc: Exception | None = None
    for attempt in range(BUILDER_RETRIES):
        try:
            code = await asyncio.to_thread(
                _call_builder_agent,
                api_key=api_key,
                slide_spec=slide_spec,
                presentation_title=presentation_title,
                section_name=section_name,
                all_sections=all_sections,
            )
            return cidx, code, None
        except anthropic.AuthenticationError:
            # Auth errors are fatal across the whole stream — don't retry.
            raise
        except (anthropic.RateLimitError, anthropic.APIStatusError) as e:
            last_exc = e
            status = getattr(e, "status_code", None)
            retriable = status in (429, 500, 502, 503, 529)
            if not retriable or attempt == BUILDER_RETRIES - 1:
                break
            await asyncio.sleep(delay)
            delay *= 2
        except Exception as e:
            last_exc = e
            if attempt == BUILDER_RETRIES - 1:
                break
            await asyncio.sleep(delay)
            delay *= 2

    msg = f"{type(last_exc).__name__ if last_exc else 'BuilderError'}: {last_exc}"
    logger.warning("[%s] builder slide %d failed after %d attempts: %s",
                   request_id, cidx, BUILDER_RETRIES, msg)
    return cidx, None, msg


@app.post("/api/generate_v2")
async def generate_v2(req: GenerateV2Request, request: Request):
    """
    V2: Hybrid generate with per-slide Builder Agent calls.
    Streams Server-Sent Events so the frontend can show per-slide progress.

    Event types:
      interpreting          — starting (no data)
      building_slide        — {current, total, title}  (incremental, parallel-aware)
      slide_error           — {slide_index, title, error}  (non-fatal)
      finalizing            — assembling .pptx (no data)
      done                  — {artifact_id, download_token, filename, warning_count}
      error                 — {message}    (fatal, stream ends)
    Heartbeat ": ping" comments are emitted every SSE_HEARTBEAT_SEC seconds to
    keep proxies (Cloudflare, corporate edges) from killing idle SSE connections.
    """
    request_id = getattr(request.state, "request_id", uuid.uuid4().hex[:12])
    if not req.slide_plan:
        raise HTTPException(400, "slide_plan cannot be empty")
    ok, reason = _validate_v2_slide_plan(req.slide_plan)
    if not ok:
        raise HTTPException(400, reason)

    slide_plan = req.slide_plan
    api_key    = req.api_key
    started    = time.monotonic()
    section_count   = len(slide_plan.get("sections", []))
    content_count   = sum(
        1 for s in slide_plan.get("sections", [])
        for sl in s.get("slides", [])
        if sl.get("slide_type", "content") == "content"
    )
    logger.info(
        "[%s] generate_v2 start sections=%d content_slides=%d concurrency=%d",
        request_id, section_count, content_count, BUILDER_CONCURRENCY,
    )

    # ── Outbound SSE queue + heartbeat ────────────────────────────────────────
    out_q: asyncio.Queue[str | None] = asyncio.Queue()

    async def emit(event: str, data: Any = None) -> None:
        await out_q.put(_sse(event, data))

    async def emit_comment(comment: str) -> None:
        # SSE comment frame — ignored by parsers but keeps the connection warm.
        await out_q.put(f": {comment}\n\n")

    async def heartbeat() -> None:
        try:
            while True:
                await asyncio.sleep(SSE_HEARTBEAT_SEC)
                await emit_comment("hb")
        except asyncio.CancelledError:
            pass

    async def disconnect_watchdog() -> None:
        # Short-circuit work if the client disconnected mid-build.
        try:
            while True:
                if await request.is_disconnected():
                    raise asyncio.CancelledError("client disconnected")
                await asyncio.sleep(2)
        except asyncio.CancelledError:
            pass

    async def producer() -> None:
        sem = asyncio.Semaphore(BUILDER_CONCURRENCY)

        try:
            await emit("interpreting")

            # Collect content slides
            all_sections   = [s["name"] for s in slide_plan.get("sections", [])]
            content_slides = []  # (section_name, slide_spec, cidx)
            for section in slide_plan.get("sections", []):
                for slide_spec in section.get("slides", []):
                    if slide_spec.get("slide_type", "content") == "content":
                        content_slides.append(
                            (section["name"], slide_spec, len(content_slides))
                        )

            total = len(content_slides)
            content_code_map: dict[int, str] = {}
            builder_failures: set[int] = set()
            completed = 0
            completed_lock = asyncio.Lock()

            async def _one(section_name: str, slide_spec: dict, cidx: int) -> None:
                nonlocal completed
                async with sem:
                    if await request.is_disconnected():
                        raise asyncio.CancelledError("client disconnected")
                    cidx_, code, err = await _builder_with_retry(
                        api_key=api_key,
                        slide_spec=slide_spec,
                        presentation_title=slide_plan.get("presentation_title", ""),
                        section_name=section_name,
                        all_sections=all_sections,
                        request_id=request_id,
                        cidx=cidx,
                    )
                    title = slide_spec.get("title", f"Slide {cidx_ + 1}")
                    async with completed_lock:
                        completed += 1
                        await emit("building_slide", {
                            "current": completed,
                            "total":   total,
                            "title":   title,
                        })
                    if err is None and code is not None:
                        content_code_map[cidx_] = code
                    else:
                        builder_failures.add(cidx_)
                        content_code_map[cidx_] = ""
                        await emit("slide_error", {
                            "slide_index": cidx_,
                            "title":       title,
                            "error":       err or "Builder Agent failed",
                        })

            # Run builder calls concurrently, bounded by Semaphore
            await asyncio.gather(*[
                _one(sn, ss, idx) for (sn, ss, idx) in content_slides
            ])

            if await request.is_disconnected():
                raise asyncio.CancelledError("client disconnected")

            # ── Assemble .pptx ─────────────────────────────────────────────
            await emit("finalizing")

            engine = InvescoreTemplateEngine(str(TEMPLATE_PATH), str(BRAND_GUIDE_PATH))
            output_path, engine_warnings = await asyncio.to_thread(
                engine.create_presentation_v2,
                slide_plan,
                content_code_map,
            )

            warning_ids = set(builder_failures)
            for warning in engine_warnings:
                warning_idx = warning.get("builder_index")
                if warning_idx in warning_ids:
                    continue
                warning_ids.add(warning_idx)
                await emit("slide_error", {
                    "slide_index": warning_idx,
                    "title": warning.get("title"),
                    "error": warning.get("message"),
                })

            artifact = await asyncio.to_thread(
                _store_download_artifact,
                output_path,
                slide_plan.get("presentation_title", "Presentation"),
            )

            await emit("done", {
                "artifact_id": artifact["artifact_id"],
                "download_token": artifact["download_token"],
                "filename": artifact["filename"],
                "warning_count": len(warning_ids),
            })
            logger.info(
                "[%s] generate_v2 done content_slides=%d fallbacks=%d duration_s=%.2f artifact=%s",
                request_id, content_count, len(warning_ids),
                time.monotonic() - started, artifact["artifact_id"],
            )

        except asyncio.CancelledError:
            logger.info(
                "[%s] generation cancelled (client disconnected) duration_s=%.2f",
                request_id, time.monotonic() - started,
            )
            # Don't emit — client is gone.
        except anthropic.AuthenticationError:
            await emit("error", {"message": "Invalid API key"})
        except anthropic.RateLimitError:
            await emit("error", {"message": "API rate limit exceeded"})
        except Exception as e:
            logger.exception("[%s] generate_v2 fatal error: %s", request_id, e)
            await emit("error", {"message": f"Generation failed (ref: {request_id})"})
        finally:
            await out_q.put(None)  # sentinel — close the stream

    async def event_stream() -> AsyncGenerator[str, None]:
        producer_task = asyncio.create_task(producer())
        heartbeat_task = asyncio.create_task(heartbeat())
        watchdog_task = asyncio.create_task(disconnect_watchdog())
        try:
            while True:
                item = await out_q.get()
                if item is None:
                    return
                yield item
        finally:
            heartbeat_task.cancel()
            watchdog_task.cancel()
            if not producer_task.done():
                producer_task.cancel()
            # Drain so background tasks finish cleanly.
            for t in (producer_task, heartbeat_task, watchdog_task):
                try:
                    await t
                except (asyncio.CancelledError, Exception):
                    pass

    return StreamingResponse(
        event_stream(),
        media_type="text/event-stream",
        headers={
            "Cache-Control":     "no-cache",
            "X-Accel-Buffering": "no",
            "Connection":        "keep-alive",
        },
    )


# ── Download endpoint (serves the assembled .pptx after SSE done) ─────────────

@app.get("/api/download/{artifact_id}")
async def download(artifact_id: str, x_download_token: str | None = Header(default=None)):
    """
    Serve a generated .pptx via an opaque artifact id + download token.

    Status codes intentionally collapsed: any failure (bad id, missing token,
    expired, wrong token, file gone) returns 404 so an attacker cannot probe
    artifact existence. Periodic cleanup is handled by the background task —
    no inline cleanup here (avoids contention with concurrent generations).
    """
    NOT_FOUND = HTTPException(404, "File not found or no longer available")

    if not re.fullmatch(r"[0-9a-f]{32}", artifact_id):
        raise NOT_FOUND
    if not x_download_token or len(x_download_token) > 256:
        raise NOT_FOUND

    meta_path = _artifact_meta_path(artifact_id)
    file_path = _artifact_file_path(artifact_id)
    if not meta_path.exists() or not file_path.exists():
        raise NOT_FOUND

    try:
        metadata = json.loads(meta_path.read_text(encoding="utf-8"))
    except Exception:
        raise NOT_FOUND

    if float(metadata.get("expires_at", 0)) <= time.time():
        # Expired — let the background sweeper actually unlink (avoid races
        # with a request still streaming the file).
        raise NOT_FOUND

    expected_hash = metadata.get("token_sha256", "")
    presented_hash = _hash_download_token(x_download_token)
    if not expected_hash or not secrets.compare_digest(expected_hash, presented_hash):
        raise NOT_FOUND

    return FileResponse(
        path=str(file_path),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename=metadata.get("filename", f"{artifact_id}.pptx"),
    )


# ── Entry point ────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    import uvicorn
    uvicorn.run("main:app", host="0.0.0.0", port=8000, reload=True)
